from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from datetime import date, datetime

import constants as const 
import data_utils as du

def create_blank_presentation(template='./template/_template.pptx'):
    prs = Presentation(template)
    return prs

def save_exit(prs, modifier=""):
    # Save the PowerPoint presentation
    today = date.today()
    current_time = datetime.now().strftime("%H%M")  # Get current hour and minute
    output_pptx = f'./output/PEA_Project_Report_{today.strftime("%y%m%d")}_{current_time}'
    output_pptx += modifier
    output_pptx += ".pptx"

    prs.save(output_pptx)

def create_title_slide(prs, title=""):
    section_slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[3])  # Blank slide layout
    # Set slide title to Objective
    title_shape = section_slide.shapes.title
    title_shape.text = f'{title}'

def set_title(slide, title_text=""):
    # Set slide title to Project Owner's name
    title_shape = slide.shapes.title
    title_shape.text = title_text

def create_project_button(slide, left, top, status, contents_text, OVERRIDE=""):
    if OVERRIDE == "":
        BUTTON_DEF = const.PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = OVERRIDE

    # Add a rounded rectangle shape
    rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            left, top, BUTTON_DEF['rectangle_width'], BUTTON_DEF['rectangle_height'])

    FILL_COLOUR = const.STATUS_COLOUR.get(status, const.ThemeColors.PINK)
    
    # Customize the rectangle Fill
    fill = rounded_rectangle.fill
    fill.solid()
    fill.fore_color.theme_color = FILL_COLOUR

    # Customize the rectangle line
    line = rounded_rectangle.line
    line.color.theme_color = FILL_COLOUR      
    line.color.brightness = -0.50 

    # Remove all Shadows
    shadow = rounded_rectangle.shadow
    shadow.inherit = False
    shadow.blur_radius = Cm(0)
    shadow.distance = Cm(0)
    shadow.angle = 0
    shadow.alpha = 0

    text_box = rounded_rectangle.text_frame
    text_box.text = contents_text

    first = 1
    for paragraph in text_box.paragraphs:
        for run in paragraph.runs:
            if first == 1:
                run.font.bold = True
                first = 0
            run.font.size = Pt(12)
    first = 1

def create_body_slide_four_cols(df, prs, type_flag='ProjectOwner', title_text="", BUTTON_OVERRIDE=""):
    # Function to take contents of df (dataframe) and output onto a 4 column grid using pre-sets from constants.py
    # Set grid parameters
    columns = 4  # Number of columns in the grid
    rows = -(-len(df) // columns)  # Calculate the number of rows needed to fit all projects

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])  # Blank slide layout
    set_title(slide, title_text)

    # Set up Constants
    SLIDE_DEF = const.FOUR_COL_SLIDE_CONSTANTS

    if BUTTON_OVERRIDE == "":
        BUTTON_DEF = const.PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = BUTTON_OVERRIDE

        # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(df.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns
        status = project['Status']

        # Calculate position for the current rectangle
        left = SLIDE_DEF['start_left'] + column * (BUTTON_DEF['rectangle_width'] + SLIDE_DEF['horizontal_spacing'])
        top = SLIDE_DEF['start_top'] + row * (BUTTON_DEF['rectangle_height'] + SLIDE_DEF['vertical_spacing'])

        # Add project details to the rectangle (based on "type_flag" for specific contents)
        if type_flag == 'ProjectOwner':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Objective: {project['Objective']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"Priority: {const.get_priority_text(project['Priority'])}"

        if type_flag == 'Objective':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"{const.get_priority_text(project['Priority'])}"

        if type_flag == 'Impact':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"{const.get_priority_text(project['Priority'])}"

        if type_flag == 'OnHold':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"Project Summary: {project['Project Summary']}"

        create_project_button(slide, left, top, status, contents_text, OVERRIDE=BUTTON_DEF)

def create_OnHold_slides(df, prs, no_section=False):
    if no_section == False:
        create_title_slide(prs, f'On Hold Projects')
    #Filter the dataframe to only the OnHold projects
    on_hold = du.filter_dataframe_by_status(df, 'On Hold')
    sorted = on_hold.sort_values(by=['Primary Owner'])
    title_text = "On-Hold Projects - " + str(len(on_hold))
    create_body_slide_four_cols(sorted, prs, type_flag='OnHold', title_text=title_text, BUTTON_OVERRIDE=const.ONHOLD_BUTTON_CONSTANTS)

def create_ProjectOwner_slides(df, prs, filter=""):
    
    # Group projects by Project Owner
    grouped = df.groupby('Primary Owner')

    # Iterate through each Project Owner and their projects
    for owner, projects in grouped:
        if not filter == "":
            if filter and filter.lower() not in owner.lower():
                continue  # Skip this owner if the filter doesn't match

        title_text = owner + " - " + str(len(projects))
        sorted_projects = projects.sort_values(by=['Objective', 'Priority'])
        create_body_slide_four_cols(sorted_projects, prs, 'ProjectOwner', title_text)

def create_Objective_slides(df, prs, filter=""):

    create_title_slide(prs, 'By Objective')

    if not filter == "":
        df.query("`Primary Owner` == @filter", inplace=True)

    # Group projects by Objective
    grouped = df.groupby('Objective')

    # Iterate through each Objective and their projects
    for objective, projects in grouped:
        
        title_text = objective  + " - " + str(len(projects))
        sorted_projects = projects.sort_values(by=['Priority'])
        create_body_slide_four_cols(sorted_projects, prs, 'Objective', title_text)

def create_Impacted_section(df, prs, no_section=False, impacted_team='Training'):
    if no_section == False:
        create_title_slide(prs, f'Projects Impacting {impacted_team}')

    # Filter projects by Impacted Team
    projects = du.filter_dataframe_by_team(df, impacted_team)

    sorted_projects = projects.sort_values(by=['Priority'])
    title_text = impacted_team  + " - " + str(len(projects))
    create_body_slide_four_cols(sorted_projects, prs, 'Impact', title_text) 

def create_Impacted_section_OLD(df, prs, no_section=False, impacted_team='Training'):
##NOT YET REFACTORED
    if no_section == False:
        create_title_slide(f'Projects Impacting {impacted_team}')

    # Filter projects by Impacted Team
    projects = du.filter_dataframe_by_team(df, impacted_team)

    # Set grid parameters
    columns = 4  # Number of columns in the grid
    rows = -(-len(df) // columns)  # Calculate the number of rows needed to fit all projects
        
    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])  # Blank slide layout

    # Set slide title to Objective
    title_shape = slide.shapes.title
    title_shape.text = impacted_team  + " - " + str(len(projects))

    # Calculate the width and height of each rectangle
    rectangle_width = Cm(7.8)
    rectangle_height = Cm(2.8)

    # Calculate the horizontal and vertical spacing between rectangles
    horizontal_spacing = Cm(0.2)
    vertical_spacing = Cm(0.2)

    # Initialize starting positions
    start_left = Cm(0.65)
    start_top = Cm(2)

    sorted_projects = projects.sort_values(by=['Priority'])

    # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(sorted_projects.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns

        # Calculate position for the current rectangle
        left = start_left + column * (rectangle_width + horizontal_spacing)
        top = start_top + row * (rectangle_height + vertical_spacing)

        # Add a rounded rectangle shape
        rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                left, top, rectangle_width, rectangle_height)

        FILL_COLOUR = const.STATUS_COLOUR.get(project['Status'], const.ThemeColors.PINK)

        # Customize the rectangle Fill
        fill = rounded_rectangle.fill
        fill.solid()
        fill.fore_color.theme_color = FILL_COLOUR

        # Customize the rectangle line
        line = rounded_rectangle.line
        line.color.theme_color = FILL_COLOUR      
        line.color.brightness = -0.50 

        # Remove all Shadows
        shadow = rounded_rectangle.shadow
        shadow.inherit = False
        shadow.blur_radius = Cm(0)
        shadow.distance = Cm(0)
        shadow.angle = 0
        shadow.alpha = 0

        # Add project details to the rectangle
        text = f"{project['Title']}\n"
        text += f"Owner: {project['Primary Owner']}\n"
        text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
        text += f"{const.get_priority_text(project['Priority'])}"

        text_box = rounded_rectangle.text_frame
        text_box.text = text

        first = 1
        for paragraph in text_box.paragraphs:
            for run in paragraph.runs:
                if first == 1:
                    run.font.bold = True
                    first = 0
                run.font.size = Pt(12)
        first = 1    