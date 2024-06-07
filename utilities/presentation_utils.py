from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import date, datetime
import pandas as pd
import math
import ast
from icecream import ic

#import PresentationToolKit.utilities.constants as const 
#import PresentationToolKit.utilities.data_utils as du

import utilities.constants as const 
import utilities.data_utils as du

def create_blank_presentation(template='./templates/_template.pptx'):
    prs = Presentation(template)
    return prs

def save_exit(prs, report_type="PEA_Project_Report", modifier="", folder=""):
    # Save the PowerPoint presentation
    today = date.today()
    current_time = datetime.now().strftime("%H%M")  # Get current hour and minute
    output_pptx = f'{today.strftime("%y%m%d")}_{current_time}_{report_type}'
    output_pptx += modifier
    output_pptx += ".pptx"

    save_to_location = folder + output_pptx
    prs.save(save_to_location)
    return output_pptx

def placeholder_identifier(slide):
    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
        print('%d, %s' % (phf.idx, phf.type))

def create_title_slide(prs, title=""):
    section_slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[3])  # Blank slide layout
    # Set slide title to Objective
    title_shape = section_slide.shapes.title
    title_shape.text = f'{title}'

def set_title(slide, title_text=""):
    # Set slide title to Project Owner's name
    title_shape = slide.shapes.title
    title_shape.text = title_text

def set_three_col_subtitle(slide):
    subtitle = const.THREE_COL_TITLES
    #hardcoded idx values that will be consistent for prs.slide_masters[1].slide_layouts[6] only
    slide.placeholders[27].text = subtitle['col1']
    slide.placeholders[29].text = subtitle['col2']
    slide.placeholders[31].text = subtitle['col3']
    return

def set_document_release_subtitle(slide, heading='new', date='08/09/2023'):
    subtitle = const.DOC_BOARD_TITLES
    #hardcoded idx values that will be consistent for prs.slide_masters[1].slide_layouts[7] only
    slide.placeholders[27].text = subtitle[heading]
    slide.placeholders[28].text = date
    return

def create_project_button(slide, left, top, status="", contents_text="CONTENT", OVERRIDE=""):
    if OVERRIDE == "":
        BUTTON_DEF = const.PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = OVERRIDE

    # Add a rounded rectangle shape
    rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            left, top, BUTTON_DEF['rectangle_width'], BUTTON_DEF['rectangle_height'])
    if status == "":
        FILL_COLOUR = BUTTON_DEF.get('fill', const.ThemeColors.PINK)
        FONT_COLOUR = BUTTON_DEF.get('font_colour', const.ThemeColors.WHITE)
        BORDER_COLOUR = BUTTON_DEF.get('border', const.ThemeColors.PINK)
        FILL_BRIGHTNESS = BUTTON_DEF.get('fill_brightness', 0)
    else:
        FILL_COLOUR = BUTTON_DEF['status_colors'].get(status, const.ThemeColors.PINK)
        BORDER_COLOUR = BUTTON_DEF.get('border', FILL_COLOUR)
        FONT_COLOUR = BUTTON_DEF.get('font_colour', 0)
        FILL_BRIGHTNESS = BUTTON_DEF.get('fill_brightness', 0)

    # Customize the rectangle Fill
    fill = rounded_rectangle.fill
    fill.solid()
    fill.fore_color.theme_color = FILL_COLOUR
    fill.fore_color.brightness = FILL_BRIGHTNESS

    # Customize the rectangle line
    line = rounded_rectangle.line
    line.color.theme_color = BORDER_COLOUR      
    line.color.brightness = -0.50 

    # Remove all Shadows
    shadow = rounded_rectangle.shadow
    shadow.inherit = False
    shadow.blur_radius = Cm(0)
    shadow.distance = Cm(0)
    shadow.angle = 0
    shadow.alpha = 0

    text_box = rounded_rectangle.text_frame
    text_box.text = contents_text
    text_box.vertical_anchor = MSO_ANCHOR.TOP

    first = 1
    for paragraph in text_box.paragraphs:
        for run in paragraph.runs:
            if FONT_COLOUR != 0:
                run.font.color.theme_color = FONT_COLOUR
            if first == 1:
                run.font.bold = True
                first = 0
            run.font.size = BUTTON_DEF['font_size']
    first = 1

def row_calculator(index, type=1):
    result = index / type
    row = math.floor(result)
    return row

def populate_column(df, slide, BUTTON_FORMAT, SLIDE_FORMAT, COLUMN_FORMAT, type_flag='ProjectOwner', col=1):
    for index, (_, project) in enumerate(df.iterrows()):
        # Calculate current row and column
        row = row_calculator(index, col)
        
        status = project['Status']

        # Calculate position for the current rectangle
        left = COLUMN_FORMAT['left']  

        if 'right' in COLUMN_FORMAT:
            right = COLUMN_FORMAT['right']
        else:
            right = left
        
        if not index % 2 == 0:
            left = right

        top = SLIDE_FORMAT['start_top'] + row * (BUTTON_FORMAT['rectangle_height'] + SLIDE_FORMAT['vertical_spacing'])

        # Add project details to the rectangle (based on "type_flag" for specific contents)
        if type_flag == 'ProjectOwner':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Objective: {project['Objective']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"Priority: {const.get_priority_text(project['Priority'])}"
            if project['Status'] == 'Blocked':
                contents_text += f"\nBlocked: {project['Closure Comments']}"

        if type_flag == 'Objective':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"{const.get_priority_text(project['Priority'])}"

        if type_flag == 'Impact':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}"
               
        if type_flag == 'OnHold':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += "Project Summary: "
            if not pd.isna(project['Project Summary']):
                contents_text += f"{project['Project Summary']}"

        create_project_button(slide, left, top, status, contents_text, OVERRIDE=BUTTON_FORMAT)

def create_body_slide_three_cols(df_1, df_2, df_3, prs, type_flag='ProjectOwner', title_text="", BUTTON_OVERRIDE=""):
    # Function to take contents of df (dataframe) and output onto a 3 column grid using pre-sets from constants.py
    # Set grid parameters
    columns = 3  # Number of columns in the grid

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[6])  # Blank slide layout
    set_title(slide, title_text)
    set_three_col_subtitle(slide)

    # Set up Constants
    SLIDE_DEF = const.THREE_COL_SLIDE_CONSTANTS

    #Hardcoding in the Column1 Button Dimension - fix in future revision
    COL1_BUTTON_DEF = const.THREE_COL_PROJECT_BUTTON_COL1_CONSTANTS

    if BUTTON_OVERRIDE == "":
        BUTTON_DEF = const.THREE_COL_PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = BUTTON_OVERRIDE

    COLUMN_1 = {
        'left': SLIDE_DEF['start_left_col1'],
        'right': SLIDE_DEF['start_left_col2']
    }
    COLUMN_2 = {
        'left': SLIDE_DEF['start_left_col3']
    }
    COLUMN_3 = {
        'left': SLIDE_DEF['start_left_col4']
    }
        # Iterate through each project and add a rounded rectangle (COLUMN 1)
    populate_column(df_1, slide, COL1_BUTTON_DEF, SLIDE_DEF, COLUMN_1, type_flag, col=2)
    populate_column(df_2, slide, BUTTON_DEF, SLIDE_DEF, COLUMN_2, type_flag)
    populate_column(df_3, slide, BUTTON_DEF, SLIDE_DEF, COLUMN_3, type_flag)

def create_body_slide_four_cols(df, prs, type_flag='ProjectOwner', title_text="", BUTTON_OVERRIDE=""):
    # Function to take contents of df (dataframe) and output onto a 4 column grid using pre-sets from constants.py
    # Set grid parameters
    columns = 4  # Number of columns in the grid

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])  # Blank slide layout
    set_title(slide, title_text)

    # Set up Constants
    SLIDE_DEF = const.FOUR_COL_SLIDE_CONSTANTS

    if BUTTON_OVERRIDE == "":
        BUTTON_DEF = const.PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = BUTTON_OVERRIDE

        # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(df.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns
        status = project['Status']

        # Calculate position for the current rectangle
        left = SLIDE_DEF['start_left'] + column * (BUTTON_DEF['rectangle_width'] + SLIDE_DEF['horizontal_spacing'])
        top = SLIDE_DEF['start_top'] + row * (BUTTON_DEF['rectangle_height'] + SLIDE_DEF['vertical_spacing'])

        # Add project details to the rectangle (based on "type_flag" for specific contents)
        if type_flag == 'ProjectOwner':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Objective: {project['Objective']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"Priority: {const.get_priority_text(project['Priority'])}"
            if project['Status'] == 'Blocked':
                contents_text += f"\nBlocked: {project['Closure Comments']}"

        if type_flag == 'Objective':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"{const.get_priority_text(project['Priority'])}"

        if type_flag == 'Impact':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"{const.get_priority_text(project['Priority'])}"

        if type_flag == 'OnHold':
            contents_text = f"{project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
            contents_text += f"Project Summary: {project['Project Summary']}"

        if type_flag == 'Release Forecast':
            contents_text = f"Document: {project['Doc Reference']}\n"
            contents_text += f"Title: {project['Title']}\n"
            contents_text += f"Owner: {project['Primary Owner']}"

        create_project_button(slide, left, top, status, contents_text, OVERRIDE=BUTTON_DEF)

def create_body_slide_four_cols_all_projects(df, prs, type_flag='ProjectOwner', title_text="", BUTTON_OVERRIDE=""):
    # Function to take contents of df (dataframe) and output onto a 4 column grid using pre-sets from constants.py
    # Set grid parameters
    columns = 4  # Number of columns in the grid

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])  # Blank slide layout
    set_title(slide, title_text)

    # Set up Constants
    SLIDE_DEF = const.FOUR_COL_SLIDE_CONSTANTS

    if BUTTON_OVERRIDE == "":
        BUTTON_DEF = const.PROJECT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = BUTTON_OVERRIDE

        # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(df.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns
        status = project['Status']

        # Calculate position for the current rectangle
        left = SLIDE_DEF['start_left'] + column * (BUTTON_DEF['rectangle_width'] + SLIDE_DEF['horizontal_spacing'])
        top = SLIDE_DEF['start_top'] + row * (BUTTON_DEF['rectangle_height'] + SLIDE_DEF['vertical_spacing'])

        # Add project details to the rectangle (based on "type_flag" for specific contents)
        if type_flag == 'ProjectOwner':
            contents_text = f"{project['Title']}"

        if type_flag == 'Objective':
            contents_text = f"{project['Title']}"

        if type_flag == 'Impact':
            contents_text = f"{project['Title']}"

        if type_flag == 'OnHold':
            contents_text = f"{project['Title']}"

        create_project_button(slide, left, top, status, contents_text, OVERRIDE=BUTTON_DEF)

def create_OnHold_slides(df, prs, no_section=False):
    if no_section == False:
        create_title_slide(prs, f'On Hold Projects')
    #Filter the dataframe to only the OnHold projects
    on_hold = du.filter_dataframe_by_status(df, 'On Hold')
    sorted = on_hold.sort_values(by=['Primary Owner'])
    title_text = "On-Hold Projects - " + str(len(on_hold))
    create_body_slide_four_cols(sorted, prs, type_flag='OnHold', title_text=title_text, BUTTON_OVERRIDE=const.ONHOLD_BUTTON_CONSTANTS)

def create_AllProjects_slide(df, prs, filter=""):
    
    title_text = "All Projects"
    create_body_slide_four_cols_all_projects(df, prs, 'ProjectOwner', title_text, const.ALL_PROJECTS_BUTTON_CONSTANTS)

def create_ProjectOwner_slides(df, prs, filter=""):
    
    # Group projects by Project Owner
    grouped = df.groupby('Primary Owner')

    # Iterate through each Project Owner and their projects
    for owner, projects in grouped:
        if not filter == "":
            if filter and filter.lower() not in owner.lower():
                continue  # Skip this owner if the filter doesn't match

        title_text = owner + " - " + str(len(projects))
        sorted_projects = projects.sort_values(by=['Priority', 'Objective'])
        create_body_slide_four_cols(sorted_projects, prs, 'ProjectOwner', title_text)

def create_Objective_slides(df, prs, filter=""):

    create_title_slide(prs, 'By Objective')

    if not filter == "":
        df = df[df['Primary Owner'].str.contains(filter, case=False, na=False)]


    # Group projects by Objective
    grouped = df.groupby('Objective')

    # Iterate through each Objective and their projects
    for objective, projects in grouped:
        
        title_text = objective  + " - " + str(len(projects))
        sorted_projects = projects.sort_values(by=['Priority'])
        create_body_slide_four_cols(sorted_projects, prs, 'Objective', title_text)

def create_Impacted_section(df, prs, no_section=False, impacted_team='Training'):
    if no_section == False:
        create_title_slide(prs, f'Projects Impacting {impacted_team}')

    # Filter projects by Impacted Team
    projects = du.filter_dataframe_by_team(df, impacted_team)
    sorted_projects = projects.sort_values(by=['Priority'])

    title_text = impacted_team  + " - " + str(len(projects))
#    create_body_slide_four_cols(sorted_projects, prs, 'Impact', title_text)

    df_triage = du.filter_dataframe_by_staging(projects, 'Triage')
    df_analysis = du.filter_dataframe_by_staging(projects, 'Analysis')
    df_alpha = du.filter_dataframe_by_staging(projects, 'Alpha Test')
    df_beta = du.filter_dataframe_by_staging(projects, 'Beta Test')
    df_rollout = du.filter_dataframe_by_staging(projects, 'Roll-out')

    col1 = du.combine_dataframe(df_triage, df_analysis)
    col2 = du.combine_dataframe(df_alpha, df_beta)
    col3 = df_rollout

    create_body_slide_three_cols(col1, col2, col3, prs, 'Impact', title_text)

def create_Impacted_section_OLD(df, prs, no_section=False, impacted_team='Training'):
##NOT YET REFACTORED
    if no_section == False:
        create_title_slide(f'Projects Impacting {impacted_team}')

    # Filter projects by Impacted Team
    projects = du.filter_dataframe_by_team(df, impacted_team)

    # Set grid parameters
    columns = 4  # Number of columns in the grid
    rows = -(-len(df) // columns)  # Calculate the number of rows needed to fit all projects
        
    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])  # Blank slide layout

    # Set slide title to Objective
    title_shape = slide.shapes.title
    title_shape.text = impacted_team  + " - " + str(len(projects))

    # Calculate the width and height of each rectangle
    rectangle_width = Cm(7.8)
    rectangle_height = Cm(2.8)

    # Calculate the horizontal and vertical spacing between rectangles
    horizontal_spacing = Cm(0.2)
    vertical_spacing = Cm(0.2)

    # Initialize starting positions
    start_left = Cm(0.65)
    start_top = Cm(2)

    sorted_projects = projects.sort_values(by=['Priority'])

    # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(sorted_projects.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns

        # Calculate position for the current rectangle
        left = start_left + column * (rectangle_width + horizontal_spacing)
        top = start_top + row * (rectangle_height + vertical_spacing)

        # Add a rounded rectangle shape
        rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                left, top, rectangle_width, rectangle_height)

        FILL_COLOUR = const.STATUS_COLOUR.get(project['Status'], const.ThemeColors.PINK)

        # Customize the rectangle Fill
        fill = rounded_rectangle.fill
        fill.solid()
        fill.fore_color.theme_color = FILL_COLOUR

        # Customize the rectangle line
        line = rounded_rectangle.line
        line.color.theme_color = FILL_COLOUR      
        line.color.brightness = -0.50 

        # Remove all Shadows
        shadow = rounded_rectangle.shadow
        shadow.inherit = False
        shadow.blur_radius = Cm(0)
        shadow.distance = Cm(0)
        shadow.angle = 0
        shadow.alpha = 0

        # Add project details to the rectangle
        text = f"{project['Title']}\n"
        text += f"Owner: {project['Primary Owner']}\n"
        text += f"Staging: {const.get_staging_text(project['Staging'])}\n"
        text += f"{const.get_priority_text(project['Priority'])}"

        text_box = rounded_rectangle.text_frame
        text_box.text = text

        first = 1
        for paragraph in text_box.paragraphs:
            for run in paragraph.runs:
                if first == 1:
                    run.font.bold = True
                    first = 0
                run.font.size = Pt(12)
        first = 1    

def create_project_section(df, prs, no_section=False):
    if no_section == False:
        create_title_slide(prs, f'Project Details')
    
    #DF Filtering Logic goes here

    for index, (_, project) in enumerate(df.iterrows()):
        create_single_project_slide(project, prs)
    return

def create_single_project_slide(project, prs, title_text=""):

    update = project['Project Updates']
    if pd.isna(update):
       update = " " 
    else:
        update = du.convert_html_to_text_with_newlines(update)

    action = project['Project Actions']
    if pd.isna(action):
       action = " " 
    else:
        action = du.convert_html_to_text_with_newlines(action)

    summary = project['Project Summary']
    if pd.isna(summary):
       summary = " " 

    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[7])  # Blank slide layout
    
    set_title(slide, project['Title'])

    #Slide subtitle
    slide.placeholders[10].text = f"Project Summary: {summary}"

    #Column Subtitles
    slide.placeholders[26].text = "Project Updates" #Left
    slide.placeholders[28].text = "Project Actions" #Right

    #Column Text
    slide.placeholders[24].text = f"{update}" #Left
    slide.placeholders[27].text = f"{action}" #Right

    return

def create_document_release_section(df, prs, filter='', internal=False):
    
    if filter:
        df = df.loc[df['Release Group'] == filter]

    grouped = df.groupby(df['Release Group'].fillna('None'))

    for date, documents in grouped:
        title_text = 'Technical Releases'

        # Filtering based on 'Doc Reference'
        new_docs = documents[documents['Doc Reference'].str.endswith('NEW')]
        update_docs = documents[~documents['Doc Reference'].str.endswith('NEW')]

        sorted_new_docs = new_docs.sort_values(by=['Doc Reference'])
        sorted_update_docs = update_docs.sort_values(by=['Doc Reference'])

        create_document_release_slide(sorted_new_docs, prs, date, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='new', full_text=True, internal=internal)
        create_document_release_slide(sorted_update_docs, prs, date, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='update', full_text=True,internal=internal)
    return

def create_document_release_section_multi_filter(df, prs, filter=[], internal=False):
    
    if filter:
        df = df[df['Release Group'].isin(filter)]   

    grouped = df.groupby(df['Release Group'].fillna('None'))

    for date, documents in grouped:
        title_text = 'Technical Releases'

        # Filtering based on 'Doc Reference'
        new_docs = documents[documents['Doc Reference'].str.endswith('NEW')]
        update_docs = documents[~documents['Doc Reference'].str.endswith('NEW')]

        sorted_new_docs = new_docs.sort_values(by=['Doc Reference'])
        sorted_update_docs = update_docs.sort_values(by=['Doc Reference'])

        create_document_release_slide(sorted_new_docs, prs, date, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='new', full_text=True, internal=internal)
        create_document_release_slide(sorted_update_docs, prs, date, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='update', full_text=True,internal=internal)
    return


def create_document_changes_section(df, prs, filter=''):
    
    if filter:
        df = df.loc[df['Release Forecast'] == filter]

    grouped = df.groupby(df['Doc Reference'].fillna('None'))

    for doc, changes in grouped:
        title_text = doc

        sorted_changes = changes.sort_values(by=['Release Forecast'])

        create_document_release_slide(sorted_changes, prs, date='', title_text=title_text, BUTTON_OVERRIDE=const.DOCUMENT_BUTTON_CONSTANTS, type_flag='changes', full_text=True)
    return

def create_document_Impacted_section(df, prs, no_section=False, impacted_team='Training', group_filter='', internal=False):
    subtitle = impacted_team

    if no_section == False:
        create_title_slide(prs, f'Documents Impacting {impacted_team}')

    if group_filter:
        subtitle = group_filter

    # Filter projects by Impacted Team
    documents = du.filter_dataframe_by_team(df, impacted_team)
    sorted_projects = documents.sort_values(by=['Doc Reference'])

    title_text = impacted_team  + " - " + str(len(documents))

    create_document_release_slide(documents, prs, subtitle, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='release_impact', full_text=True, internal=internal)

def create_document_Impacted_section_multi_filter(df, prs, no_section=False, impacted_team='Training', group_filter=[], internal=False):
    subtitle = impacted_team

    if no_section == False:
        create_title_slide(prs, f'Documents Impacting {impacted_team}')

    if group_filter:
        subtitle = " and ".join(group_filter)

    # Filter projects by Impacted Team
    documents = du.filter_dataframe_by_team(df, impacted_team)
    sorted_projects = documents.sort_values(by=['Doc Reference'])

    title_text = impacted_team  + " - " + str(len(documents))

    create_document_release_slide(documents, prs, subtitle, title_text, const.DOCUMENT_BUTTON_CONSTANTS, type_flag='release_impact', full_text=True, internal=internal)

def create_document_release_slide(df, prs, date='08/09/2023', title_text=" ", BUTTON_OVERRIDE="", type_flag='new', full_text=False, internal=False):
    # Function to take contents of df (dataframe) and output onto a 2 column grid using pre-sets from constants.py for the Document Release Board
    columns = 2
    status_text = ""

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[7])  # Blank slide layout
    set_title(slide, title_text)
    set_document_release_subtitle(slide, type_flag, date)

    # Set up Constants
    SLIDE_DEF = const.DOC_RELEASE_SLIDE_CONSTANTS

    if BUTTON_OVERRIDE == "":
        BUTTON_DEF = const.DOCUMENT_BUTTON_CONSTANTS
    else:
        BUTTON_DEF = BUTTON_OVERRIDE

        # Iterate through each project and add a rounded rectangle
    for index, (_, project) in enumerate(df.iterrows()):
        # Calculate current row and column
        row = index // columns
        column = index % columns
        status = project['Status']

        # Calculate position for the current rectangle
        left = SLIDE_DEF['start_left'] + column * (BUTTON_DEF['rectangle_width'] + SLIDE_DEF['horizontal_spacing'])
        top = SLIDE_DEF['start_top'] + row * (BUTTON_DEF['rectangle_height'] + SLIDE_DEF['vertical_spacing'])

        if full_text:
            clean_detail = du.convert_html_to_text_with_newlines(project['Release Text'])
        else:
            clean_detail = ""

        if internal:
            status_text = f"   ||   Status: {status}"

        # Add project details to the rectangle (based on "type_flag" for specific contents)
        if type_flag == 'new':
            contents_text = f"Document: {project['Doc Reference']}{status_text}\n"
            if internal:
                contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Title: {project['Title']}\n"
            contents_text += f"Summary: {clean_detail}"
            if not internal:
                status = ""

        if type_flag == 'update':
            contents_text = f"Document: {project['Doc Reference']}{status_text}\n"
            if internal:            
                contents_text += f"Owner: {project['Primary Owner']}\n"
            contents_text += f"Title: {project['Title']}\n"
            contents_text += f"Changes: {clean_detail}"
            if not internal:
                status = ""

        if type_flag == 'release_impact':
            contents_text = f"Document(s): {project['Doc Reference']}\n"
            contents_text += f"Change Title: {project['Title']}\n"
            contents_text += f"Summary of Changes: {clean_detail}"
            if not internal:
                status = ""

        if type_flag == 'changes':
            impact_token = map_impact_to_symbols(project['Impact'])

            contents_text = f"Document: {project['Doc Reference']}   ||   Impact: {impact_token}\n"
            contents_text += f"Title: {project['Title']}\n"
            contents_text += f"Detail: {clean_detail}"

        create_project_button(slide, left, top, status, contents_text, OVERRIDE=BUTTON_DEF)
    return

def map_impact_to_symbols(impact_str: str):
    try:
        impact_list = ast.literal_eval(impact_str)
    except (ValueError, SyntaxError):
        return {}
    
    impact_symbols = {}
    for impact in impact_list:
        impact_symbols[impact] = const.IMPACT_SYMBOL_MAPPING.get(impact, "UNKNOWN")
    
    return impact_symbols